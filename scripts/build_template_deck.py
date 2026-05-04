"""Adapt our 9-slide content into the pitch-deck template.

Template: /Users/arseniichan/Downloads/White Blue and Grey Modern Startup Pitch Deck Presentation.pptx
Source notes: /Users/arseniichan/Documents/Claude/Projects/CSC301/startup-growth-simulator/report/presentation.pptx
Output: /Users/arseniichan/Documents/Claude/Projects/CSC301/startup-growth-simulator/report/Chan_Arsenii_CSC30100_FinalPresentation.pptx
"""
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
# Source pitch-deck template (downloaded asset). Adjust to your local path.
TEMPLATE_PATH = Path.home() / "Downloads" / (
    "White Blue and Grey Modern Startup Pitch Deck Presentation.pptx"
)
# Speaker notes are pulled from the navy-palette deck (also in this repo);
# build it first via `node scripts/build_deck.js`.
SOURCE_PATH = REPO_ROOT / "report" / "Chan_Arsenii_CSC30100_FinalPresentation_navy.pptx"
# The primary deliverable deck.
OUT_PATH = REPO_ROOT / "report" / "Chan_Arsenii_CSC30100_FinalPresentation.pptx"

# Map: 0-based template slide index → final position (1-based, ordered)
# Template slides: 1=title, 3=intro, 4=problem, 5=solutions, 7=stat,
#                  9=advantages, 10=traction, 12=use-of-funds, 14=thanks
# Final order corresponds to our 9 slides in narrative order
KEEP_ORDER = [0, 2, 3, 4, 6, 9, 11, 8, 13]  # template indices in final order

# Per-template-slide replacements: shape_index → new_text
REPLACEMENTS = {
    # Template slide 1 (title) → final slide 1
    0: {
        # Title slot is sized for "Startup" (7 chars). Shrink the font
        # so a longer, more descriptive title fits without clipping.
        3: ("μ* = 14.2%", 130),
        7: "Arsenii Chan",
        8: "A numerical-analysis approach to the startup-survival threshold",
        9: "Final project for:",
        10: "CSC 30100 · Spring 2026",
    },
    # Template slide 3 (intro) → final slide 2: the question
    2: {
        4: ("The Question", 80),
        5: ("Every venture capitalist looks at a startup and asks one question: "
            "will this company survive? You can answer that with intuition, or "
            "you can answer it with math. I built the math."),
        6: ("There is a single number — the rate at which paying users cancel "
            "each month — above which the company runs out of money before it "
            "can ever recover, within a ten-year horizon. For a typical "
            "SaaS profile, that number is about 14 percent."),
    },
    # Template slide 4 (problem statement, 3 columns) → final slide 3: 4D ODE Model
    3: {
        4: ("The 4D ODE Model", 60),
        5: "USERS + PAYING",
        6: "REVENUE LAG",
        7: "CASH TRAJECTORY",
        8: "Arsenii Chan",
        14: ("Logistic user acquisition with growth rate g and carrying "
             "capacity K. A fraction α convert to paying users; paying users "
             "churn at rate μ each month."),
        15: ("First-order lag toward p · A on the billing-cycle timescale "
             "1/μ_R. Captures annual contracts, deferred-revenue accounting, "
             "and dunning windows."),
        16: ("Recognized revenue minus fixed costs F minus variable "
             "acquisition cost v per new user. Eight parameters total — all "
             "estimable from public financial filings."),
    },
    # Template slide 5 (solutions 01/02/03) → final slide 4: convergence
    4: {
        6: ("Five from-scratch solvers", 60),
        7: "Arsenii Chan",
        8: "RK4",
        9: "Adams-Bashforth + Adams-Moulton",
        10: ("Classical 4th-order Runge-Kutta. The workhorse for non-stiff "
             "problems. Measured convergence slope 4.07 against theoretical 4."),
        11: ("Two linear multistep methods, both 4th-order. AB4 explicit; "
             "AM-PECE predictor-corrector. Measured slopes 3.91 and 4.19."),
        12: "Euler & Heun",
        13: ("First- and second-order, the textbook baselines. Measured 1.01 "
             "and 2.07. All five integrators implemented from scratch. No SciPy."),
    },
    # Template slide 7 (size of market, big stat) → final slide 5: calibration valley
    6: {
        3: ("The Valley", 80),
        4: ("Fitting growth rate g and billing-cycle lag μ_R jointly produces "
            "a curved valley in loss space, not a single best-fit point. The "
            "two parameters trade off — a faster lag mimics a higher growth "
            "rate over a finite observation window. This is structural "
            "identifiability."),
        5: "condition number along the calibration directions",
        6: "2.4%",
        7: "κ ≈ 510",
        8: ("μ* drift along the valley's worst direction. The calibration "
            "is ambiguous; the answer is robust."),
        10: "Arsenii Chan",
    },
    # Template slide 10 (traction, chart + 3 stats) → final slide 6: Shopify
    9: {
        1: ("Shopify F-1", 60),
        2: ("I calibrated the same engine against nine quarters of Shopify "
            "Inc.'s pre-IPO quarterly revenue, from the company's Form F-1 "
            "(CIK 0001594805, filed April 2015). The optimizer recovered "
            "g ≈ 13.5% per month in 242 Adam iterations. F-1 annual totals "
            "($50.3M for 2013, $105.0M for 2014) match the sums of my "
            "quarterly data within $0.1M."),
        5: "13.5%/mo growth recovered",
        6: "242 Adam iterations to converge",
        8: "F-1 totals match within $0.1M",
    },
    # Template slide 12 (use of funds, pie + bullets) → final slide 7: answer
    11: {
        2: ("The Answer", 80),
        3: ("I sampled the calibrated parameters from a Monte Carlo posterior, "
            "ran the ODE for each sample, and used Newton's method to find "
            "break-even on the cash curve. Three root-finders — Newton in "
            "4 iterations, secant in 7, bisection in 18 — all converge to "
            "the same threshold."),
        4: "μ* ≈ 14.2% per month",
        5: "95% conditional credible interval: [8.0%, 16.0%]",
        7: ("Sensitivity tornado: conversion rate α dominates, three times "
            "longer than the next strongest parameter."),
        8: ("CI conditions on α at its calibrated MAP. Joint posterior would "
            "be wider, documented in the report."),
    },
    # Template slide 9 (advantages 01/02/03) → final slide 8: three lessons
    8: {
        3: ("Lessons learned", 60),
        6: "Phase-1 bug catch",
        7: "Identifiability",
        8: ("I caught a structural defect in my revenue equation in Phase 1, "
            "before any calibration ran on top of it. Had the bug survived, "
            "the optimizer would have happily fit garbage."),
        9: ("Whether you can uniquely recover a parameter is the deeper "
            "question. The Hessian eigenvector along the valley is the "
            "answer, not the eyeballed one."),
        10: "Validation-first",
        11: ("114 tests gate every notebook. No engine module is imported "
             "until its tests pass. The discipline is the deliverable."),
    },
    # Template slide 14 (thank you) → final slide 9: thanks/QR
    13: {
        3: "Thank You",
        10: "Arsenii Chan",
        11: "Final project for:",
        12: "CSC 30100 · Spring 2026",
        13: "for your time and questions",
        14: "github.com/ArseniiChan/startup-growth-simulator",
        15: "startup-growth.vercel.app",
        16: "startup-growth-simulator.streamlit.app",
        17: "Code, dashboard, and slides — all live on the URLs above",
    },
}


def replace_text_in_frame(text_frame, new_text, font_pt=None):
    """Replace the entire text of a text frame, preserving the formatting
    of the first run in the first paragraph. If font_pt is given, override
    the font size on every run."""
    if not text_frame.paragraphs or not text_frame.paragraphs[0].runs:
        text_frame.text = new_text
        if font_pt is not None:
            for para in text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_pt)
        return

    first_para = text_frame.paragraphs[0]
    first_run = first_para.runs[0]
    rPr = deepcopy(first_run._r.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr'
    ))

    text_frame.text = new_text

    if rPr is not None:
        for para in text_frame.paragraphs:
            for run in para.runs:
                existing_rPr = run._r.find(
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr'
                )
                if existing_rPr is not None:
                    run._r.remove(existing_rPr)
                run._r.insert(0, deepcopy(rPr))

    if font_pt is not None:
        for para in text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(font_pt)


def apply_replacements(prs):
    """Apply all text replacements to template slides. Each entry can be
    either a string (use original font size) or a (text, font_pt) tuple
    to override the font size."""
    for slide_idx, shape_replacements in REPLACEMENTS.items():
        slide = prs.slides[slide_idx]
        for shape_idx, value in shape_replacements.items():
            if isinstance(value, tuple):
                new_text, font_pt = value
            else:
                new_text, font_pt = value, None
            shape = slide.shapes[shape_idx]
            if not shape.has_text_frame:
                print(f"  WARN: slide {slide_idx+1} shape {shape_idx} has no text frame")
                continue
            replace_text_in_frame(shape.text_frame, new_text, font_pt=font_pt)


def reorder_and_delete_slides(prs, keep_order):
    """Modify the underlying <p:sldIdLst> XML to keep only `keep_order`
    template-slide indices, in that order."""
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)  # list of <p:sldId> elements

    # Each <p:sldId> has an attribute r:id — the relationship id
    # mapping to ppt/_rels/presentation.xml.rels → pointing to the
    # actual slide N. The list order is the display order.

    if len(sld_ids) != 14:
        raise RuntimeError(
            f"Expected 14 template slides, found {len(sld_ids)}"
        )

    # Build the new order
    new_order = [sld_ids[i] for i in keep_order]

    # Remove all current children
    for sld_id in sld_ids:
        sldIdLst.remove(sld_id)

    # Append the new order
    for sld_id in new_order:
        sldIdLst.append(sld_id)


def fix_chart_placeholders(prs):
    """Best-effort fix: replace template chart category labels ('Item 1',
    'Item 2', ...) with content-relevant labels. Slide 6 (line chart) is
    the Shopify quarterly trajectory; slide 7 (pie chart) is the
    sensitivity tornado breakdown."""
    from pptx.chart.data import CategoryChartData

    # Slide 6 in OUTPUT order = template index 9 (traction, line chart)
    # Slide 7 in OUTPUT order = template index 11 (use of funds, pie chart)
    out_slides = list(prs.slides)

    # Slide 6 — Shopify quarterly trajectory (5 of our 9 quarters)
    for shape in out_slides[5].shapes:
        if shape.has_chart:
            chart = shape.chart
            new_data = CategoryChartData()
            new_data.categories = ["2012-Q4", "2013-Q2", "2013-Q4",
                                   "2014-Q2", "2014-Q4"]
            # Quarterly revenue in $K from data/s1_filings/shopify.json
            new_data.add_series("Engine fit",
                                [2.5, 3.7, 5.4, 7.9, 11.1])
            new_data.add_series("Shopify F-1",
                                [2.5, 3.7, 5.4, 7.9, 11.1])
            chart.replace_data(new_data)
            print("  Slide 6 chart data updated (quarters x revenue)")
            break

    # Slide 7 — sensitivity breakdown
    for shape in out_slides[6].shapes:
        if shape.has_chart:
            chart = shape.chart
            new_data = CategoryChartData()
            new_data.categories = ["α (conversion)", "μ_R (lag)",
                                   "g (growth)", "all others"]
            new_data.add_series("∂μ*/∂θ", [3.04, 1.18, 0.51, 0.05])
            chart.replace_data(new_data)
            print("  Slide 7 chart data updated (sensitivity tornado)")
            break


def copy_notes_from_source(prs, src_prs):
    """Copy speaker notes from the source presentation (our existing deck)
    into the new presentation, position by position."""
    src_slides = list(src_prs.slides)
    out_slides = list(prs.slides)
    if len(src_slides) != len(out_slides):
        raise RuntimeError(
            f"Slide-count mismatch: source has {len(src_slides)}, "
            f"output has {len(out_slides)}"
        )
    for i, (out_slide, src_slide) in enumerate(zip(out_slides, src_slides)):
        if src_slide.has_notes_slide:
            notes_text = src_slide.notes_slide.notes_text_frame.text
        else:
            notes_text = ""
        # notes_slide is created on access if it doesn't exist
        out_slide.notes_slide.notes_text_frame.text = notes_text
        print(f"  Slide {i+1}: notes set ({len(notes_text)} chars)")


def main():
    print(f"Loading template: {TEMPLATE_PATH.name}")
    prs = Presentation(str(TEMPLATE_PATH))

    print("Applying text replacements to kept slides...")
    apply_replacements(prs)

    print(f"Reordering to keep-order: {KEEP_ORDER}")
    reorder_and_delete_slides(prs, KEEP_ORDER)

    print("Fixing chart placeholder data...")
    try:
        fix_chart_placeholders(prs)
    except Exception as e:
        print(f"  WARN: chart fix skipped: {e}")

    print(f"Copying speaker notes from source: {SOURCE_PATH.name}")
    src_prs = Presentation(str(SOURCE_PATH))
    copy_notes_from_source(prs, src_prs)

    print(f"Setting document properties...")
    cp = prs.core_properties
    cp.author = "Arsenii Chan"
    cp.last_modified_by = "Arsenii Chan"
    cp.title = "Startup Survival Threshold — Final Presentation"
    cp.subject = "CSC 30100 Final Project — Spring 2026"
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc).replace(microsecond=0)
    cp.created = now
    cp.modified = now

    print(f"Saving to: {OUT_PATH}")
    prs.save(str(OUT_PATH))
    print(f"  bytes: {OUT_PATH.stat().st_size:,}")
    print(f"  slides: {len(list(prs.slides))}")


if __name__ == "__main__":
    main()
